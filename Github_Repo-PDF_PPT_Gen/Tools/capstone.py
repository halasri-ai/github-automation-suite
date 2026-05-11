# """
# integrated_github_analyzer.py - Complete GitHub Repository Analysis Suite
# Combines: Download → Analyze → Generate PDF & PPT Reports
# """

# import os
# import sys
# import subprocess
# import zipfile
# import re
# from datetime import datetime
# from textwrap import wrap

# # AWS Bedrock & LangChain
# import boto3
# from langchain_text_splitters import RecursiveCharacterTextSplitter
# from langchain_aws import ChatBedrock

# # PDF Generation
# from reportlab.lib.pagesizes import letter
# from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
# from reportlab.lib.units import inch
# from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
# from reportlab.lib import colors
# from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY

# # PPT Generation
# from pptx import Presentation
# from pptx.util import Inches, Pt
# from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
# from pptx.dml.color import RGBColor
# from pptx.enum.shapes import MSO_SHAPE


# # =============================================================================
# # CONFIGURATION
# # =============================================================================

# AWS_CONFIG = {
      'AWS_ACCESS_KEY_ID="your_aws_access_key"'
      'AWS_SECRET_ACCESS_KEY="your_aws_secret_key"'
#     
#     'region': 'us-east-1',
#     'model_id': 'anthropic.claude-3-sonnet-20240229-v1:0'
# }

# IGNORE_DIRS = {
#     'node_modules', '.git', '__pycache__', 'dist', 'build', '.venv', 'env',
#     'venv', '.idea', '__MACOSX', 'logs', 'tmp', 'temp', '.next', 'coverage',
#     '.pytest_cache', 'htmlcov', 'site-packages'
# }

# VALID_EXTENSIONS = (
#     '.py', '.js', '.ts', '.tsx', '.jsx', '.java', '.cpp', '.c', '.h', '.hpp',
#     '.html', '.css', '.scss', '.sass', '.less', '.json', '.md', '.txt',
#     '.yaml', '.yml', '.xml', '.sql', '.sh', '.bash', '.go', '.rs', '.rb',
#     '.php', '.swift', '.kt', '.scala', '.r', '.vue', '.svelte'
# )


# # =============================================================================
# # MODULE 1: GITHUB DOWNLOADER
# # =============================================================================

# def download_github_repo(repo_url, output_dir="downloaded_repo"):
#     """Download GitHub repository as ZIP and extract it"""
#     print("\n" + "="*70)
#     print("📥 DOWNLOADING GITHUB REPOSITORY")
#     print("="*70)
    
#     if not os.path.exists(output_dir):
#         os.makedirs(output_dir)
    
#     # Prepare ZIP download URL
#     if repo_url.endswith(".git"):
#         repo_url = repo_url[:-4]
    
#     # Try main branch first, then master
#     for branch in ['main', 'master']:
#         zip_url = f"{repo_url}/archive/refs/heads/{branch}.zip"
#         zip_path = os.path.join(output_dir, "repo.zip")
        
#         print(f"⬇️  Downloading from: {zip_url}")
#         result = subprocess.run(
#             ["curl", "-L", zip_url, "-o", zip_path],
#             capture_output=True,
#             text=True
#         )
        
#         if result.returncode == 0 and os.path.exists(zip_path):
#             try:
#                 with zipfile.ZipFile(zip_path, 'r') as zip_ref:
#                     zip_ref.extractall(output_dir)
#                 print(f"✅ Repository extracted to: {output_dir}")
#                 os.remove(zip_path)
                
#                 # Find extracted folder
#                 extracted_folders = [f for f in os.listdir(output_dir) 
#                                     if os.path.isdir(os.path.join(output_dir, f))]
#                 if extracted_folders:
#                     return os.path.join(output_dir, extracted_folders[0])
#                 return output_dir
                
#             except Exception as e:
#                 print(f"⚠️  Error extracting ZIP: {e}")
#                 if os.path.exists(zip_path):
#                     os.remove(zip_path)
        
#         if branch == 'main':
#             print(f"⚠️  'main' branch not found, trying 'master'...")
    
#     print("❌ Failed to download repository")
#     return None


# # =============================================================================
# # MODULE 2: CODE ANALYZER
# # =============================================================================

# def setup_bedrock_client():
#     """Setup AWS Bedrock client"""
#     print("\n🔌 Connecting to AWS Bedrock...")
    
#     bedrock_client = boto3.client(
#         service_name="bedrock-runtime",
#         region_name=AWS_CONFIG['region'],
#         aws_access_key_id=AWS_CONFIG['access_key'],
#         aws_secret_access_key=AWS_CONFIG['secret_key']
#     )
    
#     chat = ChatBedrock(
#         model_id=AWS_CONFIG['model_id'],
#         client=bedrock_client,
#         model_kwargs={"temperature": 0.3}
#     )
    
#     print("✅ Connected to AWS Bedrock")
#     return chat


# def analyze_repository(folder_path, max_files=40, project_name=None):
#     """Analyze code repository and generate comprehensive report"""
#     print("\n" + "="*70)
#     print("🔍 ANALYZING CODE REPOSITORY")
#     print("="*70)
#     print(f"📂 Target: {folder_path}")
#     print(f"📊 Max files: {max_files}")
#     print("="*70 + "\n")
    
#     text_splitter = RecursiveCharacterTextSplitter(
#         chunk_size=3000,
#         chunk_overlap=200,
#         length_function=len
#     )
    
#     try:
#         model = setup_bedrock_client()
#     except Exception as e:
#         print(f"❌ Failed to connect to AWS Bedrock: {e}")
#         return None
    
#     file_summaries = []
#     processed_files = 0
#     skipped_files = 0
    
#     print("🔎 Scanning repository...\n")
    
#     for root, dirs, files in os.walk(folder_path):
#         dirs[:] = [d for d in dirs if d not in IGNORE_DIRS]
        
#         for file in files:
#             if not file.endswith(VALID_EXTENSIONS):
#                 continue
            
#             file_path = os.path.join(root, file)
#             relative_path = os.path.relpath(file_path, folder_path)
#             processed_files += 1
            
#             if processed_files > max_files:
#                 print(f"\n⚠️  Reached maximum file limit ({max_files} files)")
#                 break
            
#             try:
#                 with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
#                     content = f.read()
                
#                 if not content.strip() or len(content.strip()) < 50:
#                     print(f"⏭️  Skipping small/empty: {relative_path}")
#                     skipped_files += 1
#                     continue
                
#                 chunks = text_splitter.split_text(content)
#                 print(f"📄 [{processed_files}/{max_files}] {relative_path} ({len(chunks)} chunks)")
                
#                 chunk_analyses = []
#                 for i, chunk in enumerate(chunks):
#                     try:
#                         prompt = f"""Analyze this code concisely (2-3 sentences):
# Focus on: What it does, key functions/classes, and main purpose.

# CODE:
# {chunk}"""
#                         response = model.invoke(prompt)
#                         chunk_analyses.append(response.content)
#                     except Exception as e:
#                         print(f"   ⚠️  Chunk {i+1} failed: {str(e)[:60]}")
#                         continue
                
#                 if chunk_analyses:
#                     file_summaries.append({
#                         'path': relative_path,
#                         'analysis': " ".join(chunk_analyses)
#                     })
            
#             except Exception as e:
#                 print(f"❌ Error processing {relative_path}: {e}")
#                 skipped_files += 1
#                 continue
        
#         if processed_files > max_files:
#             break
    
#     print(f"\n✅ Analyzed: {len(file_summaries)} files | ⏭️  Skipped: {skipped_files}\n")
    
#     if not file_summaries:
#         print("❌ No files were successfully analyzed")
#         return None
    
#     # Generate comprehensive report
#     print("🤖 Generating comprehensive analysis report...\n")
    
#     combined_summaries = "\n\n".join([
#         f"FILE: {f['path']}\nANALYSIS: {f['analysis']}" for f in file_summaries
#     ])
    
#     # Step 1: Create condensed summary
#     print("🧠 Creating condensed project summary...")
#     condense_prompt = f"""You are a senior software architect. Summarize these {len(file_summaries)} file analyses 
# into a concise overview (1000-1500 words). Focus on: project purpose, architecture, key functionalities, tech stack.

# FILE ANALYSES:
# {combined_summaries}"""
    
#     try:
#         condensed_response = model.invoke(condense_prompt)
#         condensed_summary = condensed_response.content.strip()
#         print("✅ Condensed summary created\n")
#     except Exception as e:
#         print(f"⚠️  Condensed summary failed: {e}")
#         condensed_summary = combined_summaries[:8000]
    
#     # Step 2: Generate sections
#     sections = [
#         "PROJECT SUMMARY", "TECH STACK", "FILE/FOLDER STRUCTURE", "CORE MODULES",
#         "DATA FLOW & LOGIC", "CODE QUALITY ANALYSIS", "RED FLAGS & ISSUES",
#         "SECURITY RISKS", "PERFORMANCE RISKS", "REFACTOR SUGGESTIONS"
#     ]
    
#     section_prompts = {
#         "PROJECT SUMMARY": "Explain what the project does, its purpose, and architecture.",
#         "TECH STACK": "List and describe all technologies, frameworks, and tools used.",
#         "FILE/FOLDER STRUCTURE": "Describe directory organization and key components.",
#         "CORE MODULES": "Identify and explain main modules and their functions.",
#         "DATA FLOW & LOGIC": "Explain how data flows and components interact.",
#         "CODE QUALITY ANALYSIS": "Assess readability, modularity, and best practices.",
#         "RED FLAGS & ISSUES": "Identify code smells, scalability issues, and technical debt.",
#         "SECURITY RISKS": "Discuss vulnerabilities and security considerations.",
#         "PERFORMANCE RISKS": "Analyze efficiency, resource management, and bottlenecks.",
#         "REFACTOR SUGGESTIONS": "Provide actionable improvements and optimization tips."
#     }
    
#     report_sections = []
#     for idx, section in enumerate(sections, start=1):
#         print(f"🧩 Generating section {idx}/{len(sections)}: {section}...")
        
#         section_prompt = f"""You are a senior software architect. Based on this project summary, 
# write ONLY the following section in detailed, technical English.

# PROJECT SUMMARY CONTEXT:
# {condensed_summary}

# === {section} ===
# {section_prompts[section]}"""
        
#         try:
#             response = model.invoke(section_prompt)
#             report_sections.append(f"=== {section} ===\n{response.content.strip()}\n")
#             print(f"   ✅ Section '{section}' complete\n")
#         except Exception as e:
#             report_sections.append(f"=== {section} ===\n⚠️ Failed: {e}\n")
#             print(f"   ⚠️  Failed: {e}")
    
#     report = "\n\n".join(report_sections)
#     print("✅ All sections generated successfully!\n")
    
#     # Save text report
#     output_dir = "analysis_output"
#     os.makedirs(output_dir, exist_ok=True)
#     timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
#     report_path = os.path.join(output_dir, f"analysis_{timestamp}.txt")
    
#     with open(report_path, "w", encoding="utf-8") as f:
#         f.write("="*70 + "\n")
#         f.write("CODE REPOSITORY ANALYSIS REPORT\n")
#         f.write("="*70 + "\n")
#         f.write(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n")
#         f.write(f"Project: {project_name or os.path.basename(folder_path)}\n")
#         f.write(f"Files Analyzed: {len(file_summaries)}\n")
#         f.write(f"Target Folder: {folder_path}\n")
#         f.write("="*70 + "\n\n")
#         f.write(report)
    
#     print(f"💾 Report saved to: {report_path}\n")
    
#     return report, report_path, project_name or os.path.basename(folder_path)


# # =============================================================================
# # MODULE 3: PDF GENERATOR
# # =============================================================================

# def parse_report_sections(report_text):
#     """Parse report into structured sections"""
#     sections = {}
#     text = report_text.replace("\r", "\n")
#     pattern = r"===\s*(.*?)\s*==="
#     matches = list(re.finditer(pattern, text))
    
#     if not matches:
#         return sections
    
#     for i, match in enumerate(matches):
#         section_name = match.group(1)
#         if not section_name:
#             continue
#         start = match.end()
#         end = matches[i + 1].start() if i + 1 < len(matches) else len(text)
#         section_text = text[start:end].strip()
        
#         key = section_name.strip().lower().replace(" ", "_").replace("&", "and").replace("/", "_")
#         mapping = {
#             "project_summary": "project_summary",
#             "tech_stack": "tech_stack",
#             "file_folder_structure": "file_structure",
#             "core_modules": "core_modules",
#             "data_flow_logic": "data_flow",
#             "code_quality_analysis": "code_quality",
#             "red_flags_issues": "red_flags",
#             "security_risks": "security_risks",
#             "performance_risks": "performance_risks",
#             "refactor_suggestions": "refactor_suggestions",
#         }
#         sections[mapping.get(key, key)] = section_text
    
#     return sections


# def extract_bullet_points(text):
#     """Extract bullet points from text"""
#     points = []
#     for line in text.split("\n"):
#         line = line.strip()
#         if line:
#             if line.startswith(("-", "*", "•")):
#                 cleaned = re.sub(r"^[-*•]\s*", "", line).strip()
#                 if len(cleaned) > 10:
#                     points.append(cleaned)
#             elif len(line) > 25:
#                 points.append(line)
#     if not points or len(points) < 2:
#         sentences = re.split(r"[.!?]+", text)
#         points = [s.strip() for s in sentences if len(s.strip()) > 20]
#     return points if points else [text]


# def generate_pdf(report_text, project_name, output_folder="analysis_output"):
#     """Generate professional PDF report"""
#     print("\n" + "="*70)
#     print("📄 GENERATING PDF REPORT")
#     print("="*70)
    
#     os.makedirs(output_folder, exist_ok=True)
#     timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
#     pdf_name = f"report_{timestamp}.pdf"
#     output_path = os.path.join(output_folder, pdf_name)
    
#     doc = SimpleDocTemplate(
#         output_path,
#         pagesize=letter,
#         rightMargin=0.7*inch,
#         leftMargin=0.7*inch,
#         topMargin=1*inch,
#         bottomMargin=0.75*inch
#     )
    
#     story = []
#     styles = getSampleStyleSheet()
    
#     # Custom styles
#     title_style = ParagraphStyle(
#         "Title", fontSize=30, alignment=TA_CENTER,
#         textColor=colors.HexColor("#1a237e"), fontName="Helvetica-Bold", spaceAfter=20
#     )
#     subtitle_style = ParagraphStyle(
#         "Subtitle", fontSize=14, alignment=TA_CENTER,
#         textColor=colors.HexColor("#424242"), fontName="Helvetica-Bold", spaceAfter=25
#     )
#     heading_style = ParagraphStyle(
#         "Heading", fontSize=18, textColor=colors.HexColor("#283593"),
#         backColor=colors.HexColor("#e8eaf6"), fontName="Helvetica-Bold",
#         spaceBefore=20, spaceAfter=12
#     )
#     body_style = ParagraphStyle(
#         "Body", fontSize=11, leading=16, alignment=TA_JUSTIFY,
#         textColor=colors.HexColor("#212121")
#     )
#     bullet_style = ParagraphStyle(
#         "Bullet", fontSize=11, leading=14, leftIndent=20, spaceAfter=6
#     )
    
#     # Title Page
#     story.append(Paragraph("📊 Code Analysis Report", title_style))
#     story.append(Paragraph(f"<b>{project_name}</b>", subtitle_style))
#     story.append(Paragraph(f"Generated: {datetime.now().strftime('%B %d, %Y, %I:%M %p')}", 
#                           ParagraphStyle("Meta", fontSize=9, alignment=TA_CENTER)))
#     story.append(Spacer(1, 0.3*inch))
    
#     # Parse sections
#     sections = parse_report_sections(report_text)
    
#     section_titles = {
#         "project_summary": "📋 Project Summary",
#         "tech_stack": "🛠️ Technology Stack",
#         "file_structure": "📁 File & Folder Structure",
#         "core_modules": "⚙️ Core Modules",
#         "data_flow": "🔄 Data Flow & Logic",
#         "code_quality": "✅ Code Quality Analysis",
#         "red_flags": "⚠️ Red Flags & Issues",
#         "security_risks": "🔒 Security Risks",
#         "performance_risks": "⚡ Performance Risks",
#         "refactor_suggestions": "💡 Refactor Suggestions"
#     }
    
#     for key, title in section_titles.items():
#         if key in sections and sections[key].strip():
#             story.append(Paragraph(title, heading_style))
#             story.append(Spacer(1, 0.1*inch))
            
#             content = sections[key]
            
#             if key == "tech_stack":
#                 items = extract_bullet_points(content)
#                 for item in items[:20]:
#                     story.append(Paragraph(f"• {item}", bullet_style))
#             else:
#                 points = extract_bullet_points(content)
#                 if len(points) > 1:
#                     for p in points[:15]:
#                         story.append(Paragraph(f"• {p}", bullet_style))
#                 else:
#                     for para in content.split("\n\n"):
#                         if para.strip():
#                             story.append(Paragraph(para.strip(), body_style))
            
#             story.append(Spacer(1, 0.15*inch))
    
#     doc.build(story)
#     print(f"✅ PDF created: {output_path}\n")
#     return output_path


# # =============================================================================
# # MODULE 4: PPT GENERATOR
# # =============================================================================

# def add_ppt_title_slide(prs, project_name):
#     """Add title slide to presentation"""
#     slide = prs.slides.add_slide(prs.slide_layouts[6])
    
#     bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5))
#     bg.fill.solid()
#     bg.fill.fore_color.rgb = RGBColor(26, 35, 126)
#     bg.line.fill.background()
    
#     title_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
#     title_frame = title_box.text_frame
#     p = title_frame.paragraphs[0]
#     p.text = "Code Analysis Report"
#     p.font.size = Pt(44)
#     p.font.bold = True
#     p.font.color.rgb = RGBColor(255, 255, 255)
#     p.alignment = PP_ALIGN.CENTER
    
#     project_box = slide.shapes.add_textbox(Inches(2), Inches(4), Inches(6), Inches(0.5))
#     p = project_box.text_frame.paragraphs[0]
#     p.text = project_name
#     p.font.size = Pt(24)
#     p.font.color.rgb = RGBColor(200, 230, 255)
#     p.alignment = PP_ALIGN.CENTER


# def add_ppt_content_slide(prs, title, content_points):
#     """Add content slide to presentation"""
#     slide = prs.slides.add_slide(prs.slide_layouts[6])
    
#     title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.6))
#     p = title_box.text_frame.paragraphs[0]
#     p.text = title
#     p.font.size = Pt(28)
#     p.font.bold = True
#     p.font.color.rgb = RGBColor(40, 53, 147)
    
#     content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5.5))
#     text_frame = content_box.text_frame
#     text_frame.word_wrap = True
    
#     if isinstance(content_points, list):
#         for i, point in enumerate(content_points[:8]):
#             if i == 0:
#                 p = text_frame.paragraphs[0]
#             else:
#                 p = text_frame.add_paragraph()
#             p.text = f"• {point[:300]}"
#             p.font.size = Pt(14)
#             p.space_after = Pt(8)


# def generate_ppt(report_text, project_name, output_folder="analysis_output"):
#     """Generate professional PowerPoint presentation"""
#     print("\n" + "="*70)
#     print("📊 GENERATING POWERPOINT PRESENTATION")
#     print("="*70)
    
#     os.makedirs(output_folder, exist_ok=True)
#     timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
#     ppt_name = f"presentation_{timestamp}.pptx"
#     output_path = os.path.join(output_folder, ppt_name)
    
#     prs = Presentation()
#     prs.slide_width = Inches(10)
#     prs.slide_height = Inches(7.5)
    
#     add_ppt_title_slide(prs, project_name)
    
#     sections = parse_report_sections(report_text)
    
#     section_titles = {
#         "project_summary": "📋 Project Summary",
#         "tech_stack": "🛠️ Technology Stack",
#         "file_structure": "📁 File Structure",
#         "core_modules": "⚙️ Core Modules",
#         "data_flow": "🔄 Data Flow & Logic",
#         "code_quality": "✅ Code Quality",
#         "red_flags": "⚠️ Red Flags",
#         "security_risks": "🔒 Security Risks",
#         "performance_risks": "⚡ Performance",
#         "refactor_suggestions": "💡 Suggestions"
#     }
    
#     for key, title in section_titles.items():
#         if key in sections and sections[key].strip():
#             points = extract_bullet_points(sections[key])
            
#             # Split into multiple slides if needed
#             for i in range(0, len(points), 8):
#                 slide_title = title if i == 0 else f"{title} (Part {i//8 + 1})"
#                 add_ppt_content_slide(prs, slide_title, points[i:i+8])
    
#     prs.save(output_path)
#     print(f"✅ PowerPoint created: {output_path}\n")
#     return output_path


# # =============================================================================
# # MAIN ORCHESTRATOR
# # =============================================================================

# def main():
#     """Main orchestrator function"""
#     print("\n" + "="*70)
#     print("🚀 INTEGRATED GITHUB REPOSITORY ANALYZER")
#     print("="*70)
#     print("Features: Download → Analyze → Generate PDF & PPT")
#     print("="*70 + "\n")
    
#     # Step 1: Get repository URL or folder path
#     choice = input("Enter '1' for GitHub URL or '2' for local folder: ").strip()
    
#     if choice == '1':
#         repo_url = input("📎 Enter GitHub repository URL: ").strip()
#         if not repo_url:
#             print("❌ No URL provided")
#             return
        
#         folder_path = download_github_repo(repo_url)
#         if not folder_path:
#             print("❌ Failed to download repository")
#             return
        
#         project_name = repo_url.split('/')[-1].replace('.git', '')
    
#     elif choice == '2':
#         folder_path = input("📂 Enter local folder path: ").strip()
#         if not folder_path or not os.path.exists(folder_path):
#             print("❌ Invalid folder path")
#             return
#         project_name = os.path.basename(os.path.abspath(folder_path))
    
#     else:
#         print("❌ Invalid choice")
#         return
    
#     # Step 2: Analyze repository
#     max_files = input(f"📊 Max files to analyze (default=40): ").strip()
#     max_files = int(max_files) if max_files.isdigit() else 40
    
#     result = analyze_repository(folder_path, max_files, project_name)
#     if not result:
#         print("❌ Analysis failed")
#         return
    
#     report_text, report_path, project_name = result
    
#     # Step 3: Generate documents
#     gen_docs = input("\n📄 Generate PDF and PPT? (Y/n): ").strip().lower()
#     if gen_docs != 'n':
#         pdf_path = generate_pdf(report_text, project_name)
#         ppt_path = generate_ppt(report_text, project_name)
        
#         print("\n" + "="*70)
#         print("🎉 SUCCESS! ALL DOCUMENTS GENERATED")
#         print("="*70)
#         print(f"📄 Text Report: {report_path}")
#         print(f"📄 PDF Report: {pdf_path}")
#         print(f"📊 PowerPoint: {ppt_path}")
#         print("="*70 + "\n")
#     else:
#         print(f"\n✅ Analysis complete! Report saved to: {report_path}\n")


# if __name__ == "__main__":
#     try:
#         main()
#     except KeyboardInterrupt:
#         print("\n\n⚠️  Process interrupted by user")
#     except Exception as e:
#         print(f"\n❌ Error: {e}")
#         import traceback
#         traceback.print_exc()

"""
integrated_github_analyzer.py - Complete GitHub Repository Analysis Suite
Combines: Download → Analyze → Generate PDF & PPT Reports
Enhanced with improved PDF/PPT generation from original files
"""

import os
import sys
import subprocess
import zipfile
import re
from datetime import datetime
from textwrap import wrap

# AWS Bedrock & LangChain
import boto3
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_aws import ChatBedrock

# PDF Generation
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY

# PPT Generation
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


# =============================================================================
# CONFIGURATION
# =============================================================================

AWS_CONFIG = {
    'access_key': "your_aws_access_key",
    'secret_key': "your_aws_secret_key",
    'region': 'us-east-1',
    'model_id': 'anthropic.claude-3-sonnet-20240229-v1:0'
}

IGNORE_DIRS = {
    'node_modules', '.git', '__pycache__', 'dist', 'build', '.venv', 'env',
    'venv', '.idea', '__MACOSX', 'logs', 'tmp', 'temp', '.next', 'coverage',
    '.pytest_cache', 'htmlcov', 'site-packages'
}

VALID_EXTENSIONS = (
    '.py', '.js', '.ts', '.tsx', '.jsx', '.java', '.cpp', '.c', '.h', '.hpp',
    '.html', '.css', '.scss', '.sass', '.less', '.json', '.md', '.txt',
    '.yaml', '.yml', '.xml', '.sql', '.sh', '.bash', '.go', '.rs', '.rb',
    '.php', '.swift', '.kt', '.scala', '.r', '.vue', '.svelte'
)


# =============================================================================
# MODULE 1: GITHUB DOWNLOADER
# =============================================================================

def download_github_repo(repo_url, output_dir="downloaded_repo"):
    """Download GitHub repository as ZIP and extract it"""
    print("\n" + "="*70)
    print("📥 DOWNLOADING GITHUB REPOSITORY")
    print("="*70)
    
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
    
    # Prepare ZIP download URL
    if repo_url.endswith(".git"):
        repo_url = repo_url[:-4]
    
    # Try main branch first, then master
    for branch in ['main', 'master']:
        zip_url = f"{repo_url}/archive/refs/heads/{branch}.zip"
        zip_path = os.path.join(output_dir, "repo.zip")
        
        print(f"⬇️  Downloading from: {zip_url}")
        result = subprocess.run(
            ["curl", "-L", zip_url, "-o", zip_path],
            capture_output=True,
            text=True
        )
        
        if result.returncode == 0 and os.path.exists(zip_path):
            try:
                with zipfile.ZipFile(zip_path, 'r') as zip_ref:
                    zip_ref.extractall(output_dir)
                print(f"✅ Repository extracted to: {output_dir}")
                os.remove(zip_path)
                
                # Find extracted folder
                extracted_folders = [f for f in os.listdir(output_dir) 
                                    if os.path.isdir(os.path.join(output_dir, f))]
                if extracted_folders:
                    return os.path.join(output_dir, extracted_folders[0])
                return output_dir
                
            except Exception as e:
                print(f"⚠️  Error extracting ZIP: {e}")
                if os.path.exists(zip_path):
                    os.remove(zip_path)
        
        if branch == 'main':
            print(f"⚠️  'main' branch not found, trying 'master'...")
    
    print("❌ Failed to download repository")
    return None


# =============================================================================
# MODULE 2: CODE ANALYZER
# =============================================================================

def setup_bedrock_client():
    """Setup AWS Bedrock client"""
    print("\n🔌 Connecting to AWS Bedrock...")
    
    bedrock_client = boto3.client(
        service_name="bedrock-runtime",
        region_name=AWS_CONFIG['region'],
        aws_access_key_id=AWS_CONFIG['access_key'],
        aws_secret_access_key=AWS_CONFIG['secret_key']
    )
    
    chat = ChatBedrock(
        model_id=AWS_CONFIG['model_id'],
        client=bedrock_client,
        model_kwargs={"temperature": 0.3}
    )
    
    print("✅ Connected to AWS Bedrock")
    return chat


def analyze_repository(folder_path, max_files=40, project_name=None):
    """Analyze code repository and generate comprehensive report"""
    print("\n" + "="*70)
    print("🔍 ANALYZING CODE REPOSITORY")
    print("="*70)
    print(f"📂 Target: {folder_path}")
    print(f"📊 Max files: {max_files}")
    print("="*70 + "\n")
    
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=3000,
        chunk_overlap=200,
        length_function=len
    )
    
    try:
        model = setup_bedrock_client()
    except Exception as e:
        print(f"❌ Failed to connect to AWS Bedrock: {e}")
        return None
    
    file_summaries = []
    processed_files = 0
    skipped_files = 0
    
    print("🔎 Scanning repository...\n")
    
    for root, dirs, files in os.walk(folder_path):
        dirs[:] = [d for d in dirs if d not in IGNORE_DIRS]
        
        for file in files:
            if not file.endswith(VALID_EXTENSIONS):
                continue
            
            file_path = os.path.join(root, file)
            relative_path = os.path.relpath(file_path, folder_path)
            processed_files += 1
            
            if processed_files > max_files:
                print(f"\n⚠️  Reached maximum file limit ({max_files} files)")
                break
            
            try:
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    content = f.read()
                
                if not content.strip() or len(content.strip()) < 50:
                    print(f"⏭️  Skipping small/empty: {relative_path}")
                    skipped_files += 1
                    continue
                
                chunks = text_splitter.split_text(content)
                print(f"📄 [{processed_files}/{max_files}] {relative_path} ({len(chunks)} chunks)")
                
                chunk_analyses = []
                for i, chunk in enumerate(chunks):
                    try:
                        prompt = f"""Analyze this code concisely (2-3 sentences):
Focus on: What it does, key functions/classes, and main purpose.

CODE:
{chunk}"""
                        response = model.invoke(prompt)
                        chunk_analyses.append(response.content)
                    except Exception as e:
                        print(f"   ⚠️  Chunk {i+1} failed: {str(e)[:60]}")
                        continue
                
                if chunk_analyses:
                    file_summaries.append({
                        'path': relative_path,
                        'analysis': " ".join(chunk_analyses)
                    })
            
            except Exception as e:
                print(f"❌ Error processing {relative_path}: {e}")
                skipped_files += 1
                continue
        
        if processed_files > max_files:
            break
    
    print(f"\n✅ Analyzed: {len(file_summaries)} files | ⏭️  Skipped: {skipped_files}\n")
    
    if not file_summaries:
        print("❌ No files were successfully analyzed")
        return None
    
    # Generate comprehensive report
    print("🤖 Generating comprehensive analysis report...\n")
    
    combined_summaries = "\n\n".join([
        f"FILE: {f['path']}\nANALYSIS: {f['analysis']}" for f in file_summaries
    ])
    
    # Step 1: Create condensed summary
    print("🧠 Creating condensed project summary...")
    condense_prompt = f"""You are a senior software architect. Summarize these {len(file_summaries)} file analyses 
into a concise overview (1000-1500 words). Focus on: project purpose, architecture, key functionalities, tech stack.

FILE ANALYSES:
{combined_summaries}"""
    
    try:
        condensed_response = model.invoke(condense_prompt)
        condensed_summary = condensed_response.content.strip()
        print("✅ Condensed summary created\n")
    except Exception as e:
        print(f"⚠️  Condensed summary failed: {e}")
        condensed_summary = combined_summaries[:8000]
    
    # Step 2: Generate sections
    sections = [
        "PROJECT SUMMARY", "TECH STACK", "FILE/FOLDER STRUCTURE", "CORE MODULES",
        "DATA FLOW & LOGIC", "CODE QUALITY ANALYSIS", "RED FLAGS & ISSUES",
        "SECURITY RISKS", "PERFORMANCE RISKS", "REFACTOR SUGGESTIONS"
    ]
    
    section_prompts = {
        "PROJECT SUMMARY": "Explain what the project does, its purpose, and architecture.",
        "TECH STACK": "List and describe all technologies, frameworks, and tools used.",
        "FILE/FOLDER STRUCTURE": "Describe directory organization and key components.",
        "CORE MODULES": "Identify and explain main modules and their functions.",
        "DATA FLOW & LOGIC": "Explain how data flows and components interact.",
        "CODE QUALITY ANALYSIS": "Assess readability, modularity, and best practices.",
        "RED FLAGS & ISSUES": "Identify code smells, scalability issues, and technical debt.",
        "SECURITY RISKS": "Discuss vulnerabilities and security considerations.",
        "PERFORMANCE RISKS": "Analyze efficiency, resource management, and bottlenecks.",
        "REFACTOR SUGGESTIONS": "Provide actionable improvements and optimization tips."
    }
    
    report_sections = []
    for idx, section in enumerate(sections, start=1):
        print(f"🧩 Generating section {idx}/{len(sections)}: {section}...")
        
        section_prompt = f"""You are a senior software architect. Based on this project summary, 
write ONLY the following section in detailed, technical English.

PROJECT SUMMARY CONTEXT:
{condensed_summary}

=== {section} ===
{section_prompts[section]}"""
        
        try:
            response = model.invoke(section_prompt)
            report_sections.append(f"=== {section} ===\n{response.content.strip()}\n")
            print(f"   ✅ Section '{section}' complete\n")
        except Exception as e:
            report_sections.append(f"=== {section} ===\n⚠️ Failed: {e}\n")
            print(f"   ⚠️  Failed: {e}")
    
    report = "\n\n".join(report_sections)
    print("✅ All sections generated successfully!\n")
    
    # Save text report
    output_dir = "analysis_output"
    os.makedirs(output_dir, exist_ok=True)
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    report_path = os.path.join(output_dir, f"analysis_{timestamp}.txt")
    
    with open(report_path, "w", encoding="utf-8") as f:
        f.write("="*70 + "\n")
        f.write("CODE REPOSITORY ANALYSIS REPORT\n")
        f.write("="*70 + "\n")
        f.write(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n")
        f.write(f"Project: {project_name or os.path.basename(folder_path)}\n")
        f.write(f"Files Analyzed: {len(file_summaries)}\n")
        f.write(f"Target Folder: {folder_path}\n")
        f.write("="*70 + "\n\n")
        f.write(report)
    
    print(f"💾 Report saved to: {report_path}\n")
    
    return report, report_path, project_name or os.path.basename(folder_path)


# =============================================================================
# MODULE 3: ENHANCED PDF GENERATOR (from gen_pdf.py)
# =============================================================================

def parse_analysis_report(report_text):
    """Parse report into structured sections with enhanced pattern matching"""
    sections = {}
    text = report_text.replace("\r", "\n")
    pattern = r"(?:===\s*(.*?)\s*===|([A-Z\s/&]+):)"
    matches = list(re.finditer(pattern, text))
    
    if not matches:
        print("⚠️  No recognizable section headers found.")
        return sections

    for i, match in enumerate(matches):
        section_name = match.group(1) or match.group(2)
        if not section_name:
            continue
        start = match.end()
        end = matches[i + 1].start() if i + 1 < len(matches) else len(text)
        section_text = text[start:end].strip()

        key = section_name.strip().lower().replace(" ", "_").replace("&", "and").replace("/", "_")
        mapping = {
            "project_summary": "project_summary",
            "tech_stack": "tech_stack",
            "file_folder_structure": "file_structure",
            "file_structure": "file_structure",
            "core_modules": "core_modules",
            "data_flow_and_logic": "data_flow",
            "data_flow_logic": "data_flow",
            "code_quality_analysis": "code_quality",
            "red_flags_and_issues": "red_flags",
            "red_flags_issues": "red_flags",
            "security_risks": "security_risks",
            "performance_risks": "performance_risks",
            "refactor_suggestions": "refactor_suggestions",
        }
        sections[mapping.get(key, key)] = section_text
    
    print(f"✅ Parsed {len(sections)} sections: {list(sections.keys())}")
    return sections


def extract_bullet_points(text):
    """Extract bullet points from text"""
    points = []
    for line in text.split("\n"):
        line = line.strip()
        if line:
            if line.startswith(("-", "*", "•")):
                cleaned = re.sub(r"^[-*•]\s*", "", line).strip()
                if len(cleaned) > 10:
                    points.append(cleaned)
            elif len(line) > 25:
                points.append(line)
    if not points or len(points) < 2:
        sentences = re.split(r"[.!?]+", text)
        points = [s.strip() for s in sentences if len(s.strip()) > 20]
    return points if points else [text]


def generate_pdf(report_text, project_name, output_folder="analysis_output"):
    """Generate professional PDF report with enhanced formatting"""
    print("\n" + "="*70)
    print("📄 GENERATING PDF REPORT")
    print("="*70)
    
    os.makedirs(output_folder, exist_ok=True)
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    pdf_name = f"report_{timestamp}.pdf"
    output_path = os.path.join(output_folder, pdf_name)

    # Remove any existing file
    if os.path.exists(output_path):
        os.remove(output_path)

    doc = SimpleDocTemplate(
        output_path,
        pagesize=letter,
        rightMargin=0.7*inch,
        leftMargin=0.7*inch,
        topMargin=1*inch,
        bottomMargin=0.75*inch,
        title=f"Code Analysis Report - {project_name}",
        author="Enhanced AI Code Analyzer"
    )

    story = []
    styles = getSampleStyleSheet()

    # Enhanced styles
    title_style = ParagraphStyle(
        "Title", fontSize=30, alignment=TA_CENTER,
        textColor=colors.HexColor("#1a237e"), fontName="Helvetica-Bold", spaceAfter=20
    )
    subtitle_style = ParagraphStyle(
        "Subtitle", fontSize=14, alignment=TA_CENTER,
        textColor=colors.HexColor("#424242"), fontName="Helvetica-Bold", spaceAfter=25
    )
    heading_style = ParagraphStyle(
        "Heading", fontSize=18, textColor=colors.HexColor("#283593"),
        backColor=colors.HexColor("#e8eaf6"), fontName="Helvetica-Bold",
        spaceBefore=20, spaceAfter=12
    )
    body_style = ParagraphStyle(
        "Body", fontSize=11, leading=16, alignment=TA_JUSTIFY,
        textColor=colors.HexColor("#212121")
    )
    bullet_style = ParagraphStyle(
        "Bullet", fontSize=11, leading=14, leftIndent=20, spaceAfter=6,
        textColor=colors.HexColor("#212121")
    )
    meta_style = ParagraphStyle(
        "Meta", fontSize=9, alignment=TA_CENTER, textColor=colors.HexColor("#757575")
    )

    # Title Page
    story.append(Paragraph("📊 Code Analysis Report", title_style))
    story.append(Paragraph(f"<b>{project_name}</b>", subtitle_style))
    story.append(Paragraph(f"Generated: {datetime.now().strftime('%B %d, %Y, %I:%M %p')}", meta_style))
    story.append(Paragraph("AI-Powered Repository Analyzer", meta_style))
    story.append(Spacer(1, 0.3*inch))
    story.append(Table([[""]], colWidths=[6.5*inch],
                       style=[("LINEABOVE", (0, 0), (-1, 0), 2, colors.HexColor("#3f51b5"))]))
    story.append(Spacer(1, 0.4*inch))

    print("🔍 Parsing report...")
    sections = parse_analysis_report(report_text)
    
    if not sections:
        story.append(Paragraph("⚠️ Could not parse any sections.", body_style))

    section_titles = {
        "project_summary": "📋 Project Summary",
        "tech_stack": "🛠️ Technology Stack",
        "file_structure": "📁 File & Folder Structure",
        "core_modules": "⚙️ Core Modules",
        "data_flow": "🔄 Data Flow & Logic",
        "code_quality": "✅ Code Quality Analysis",
        "red_flags": "⚠️ Red Flags & Issues",
        "security_risks": "🔒 Security Risks",
        "performance_risks": "⚡ Performance Risks",
        "refactor_suggestions": "💡 Refactor Suggestions"
    }

    for section_key, title in section_titles.items():
        if section_key not in sections or not sections[section_key].strip():
            continue

        story.append(Paragraph(title, heading_style))
        story.append(Spacer(1, 0.1*inch))
        content = sections[section_key]

        # Tech Stack - Create enhanced table
        if section_key == "tech_stack":
            items = extract_bullet_points(content)
            table_data = [["Technology / Framework / Library"]]
            for item in items[:20]:
                wrapped = "<br/>".join(wrap(item, 100))
                table_data.append([Paragraph(f"• {wrapped}", body_style)])
            tech_table = Table(table_data, colWidths=[6.3*inch])
            tech_table.setStyle(TableStyle([
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#e8eaf6")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.black),
                ("ALIGN", (0, 0), (-1, -1), "LEFT"),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("FONTNAME", (0, 0), (-1, -1), "Helvetica"),
                ("FONTSIZE", (0, 0), (-1, -1), 10),
                ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
            ]))
            story.append(tech_table)
            story.append(Spacer(1, 0.2*inch))
            continue

        # Red Flags / Security Risks - Enhanced warning boxes
        if section_key in ["red_flags", "security_risks"]:
            bg_color = colors.HexColor("#fff3e0") if section_key == "red_flags" else colors.HexColor("#ffebee")
            wrapped_lines = wrap(content, 130)
            for chunk in range(0, len(wrapped_lines), 25):
                part = " ".join(wrapped_lines[chunk:chunk + 25])
                story.append(Spacer(1, 0.05*inch))
                story.append(Paragraph(
                    f'<para backColor="{bg_color}"><font color="black">{part}</font></para>',
                    body_style))
                story.append(Spacer(1, 0.1*inch))
            continue

        # Normal Sections
        points = extract_bullet_points(content)
        if len(points) > 1:
            for p in points[:15]:
                story.append(Paragraph(f"• {' '.join(wrap(p, 120))}", bullet_style))
        else:
            for para in content.split("\n\n"):
                para = para.strip().replace("\n", " ")
                if para:
                    story.append(Paragraph(para, body_style))
        story.append(Spacer(1, 0.15*inch))

    # Footer
    story.append(Spacer(1, 0.5*inch))
    story.append(Table([[""]], colWidths=[6.4*inch],
                       style=[("LINEABOVE", (0, 0), (-1, 0), 1, colors.grey)]))
    story.append(Spacer(1, 0.1*inch))
    story.append(Paragraph("End of Analysis Report", meta_style))

    doc.build(story)

    # Verify file
    if os.path.exists(output_path) and os.path.getsize(output_path) > 1000:
        print(f"✅ PDF created successfully: {os.path.abspath(output_path)}\n")
    else:
        print("⚠️ File generated, but might not be a valid PDF.\n")
    
    return output_path


# =============================================================================
# MODULE 4: ENHANCED PPT GENERATOR (from gen_ppt.py)
# =============================================================================

def add_icon_shape(slide, emoji, left, top, size=1.2):
    """Add decorative icon to slide"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top),
        Inches(size), Inches(size)
    )

    color_map = {
        '📋': RGBColor(26, 35, 126),
        '🛠️': RGBColor(13, 71, 161),
        '📁': RGBColor(27, 94, 32),
        '⚙️': RGBColor(49, 27, 146),
        '🔄': RGBColor(1, 87, 155),
        '✅': RGBColor(27, 94, 32),
        '⚠️': RGBColor(230, 81, 0),
        '🔒': RGBColor(198, 40, 40),
        '⚡': RGBColor(245, 124, 0),
        '💡': RGBColor(123, 31, 162),
    }
    bg_color = color_map.get(emoji, RGBColor(26, 35, 126))

    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = bg_color

    text_frame = shape.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = emoji
    p.font.size = Pt(int(size * 30))
    p.alignment = PP_ALIGN.CENTER
    return shape


def add_decorative_element(slide, section_type):
    """Add decorative accent bar"""
    accent_colors = {
        'project_summary': RGBColor(26, 35, 126),
        'tech_stack': RGBColor(13, 71, 161),
        'file_structure': RGBColor(27, 94, 32),
        'core_modules': RGBColor(49, 27, 146),
        'data_flow': RGBColor(1, 87, 155),
        'code_quality': RGBColor(27, 94, 32),
        'red_flags': RGBColor(230, 81, 0),
        'security_risks': RGBColor(198, 40, 40),
        'performance_risks': RGBColor(245, 124, 0),
        'refactor_suggestions': RGBColor(123, 31, 162),
    }
    color = accent_colors.get(section_type, RGBColor(26, 35, 126))

    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.3), Inches(1.3),
        Inches(0.15), Inches(5.5)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def add_ppt_title_slide(prs, project_name):
    """Add enhanced title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(26, 35, 126)
    bg.line.fill.background()

    accent = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(2), Inches(7), Inches(3.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(40, 53, 147)
    accent.line.fill.background()

    icon = slide.shapes.add_textbox(Inches(4.2), Inches(2.3), Inches(1.6), Inches(0.8))
    p = icon.text_frame.paragraphs[0]
    p.text = "📊"
    p.font.size = Pt(72)
    p.alignment = PP_ALIGN.CENTER

    title_box = slide.shapes.add_textbox(Inches(2), Inches(3.2), Inches(6), Inches(0.7))
    p = title_box.text_frame.paragraphs[0]
    p.text = "Code Analysis Report"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    project_box = slide.shapes.add_textbox(Inches(2), Inches(4.1), Inches(6), Inches(0.5))
    p = project_box.text_frame.paragraphs[0]
    p.text = project_name
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(200, 230, 255)
    p.alignment = PP_ALIGN.CENTER

    meta_box = slide.shapes.add_textbox(Inches(2), Inches(5.8), Inches(6), Inches(0.8))
    p1 = meta_box.text_frame.paragraphs[0]
    p1.text = f"Generated: {datetime.now().strftime('%B %d, %Y')}"
    p1.font.size = Pt(14)
    p1.font.color.rgb = RGBColor(150, 180, 255)
    p1.alignment = PP_ALIGN.CENTER
    
    p2 = meta_box.text_frame.add_paragraph()
    p2.text = "AI-Powered Code Repository Analyzer"
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(150, 180, 255)
    p2.alignment = PP_ALIGN.CENTER


def add_content_slide_with_image(prs, title, emoji, content_points, section_type="bullet", part_label=None):
    """Add enhanced content slide with icons and formatting"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_decorative_element(slide, section_type)
    add_icon_shape(slide, emoji, 0.7, 1.5, 1.0)

    title_box = slide.shapes.add_textbox(Inches(2.0), Inches(1.3), Inches(7.5), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = f"{title}{(' — ' + part_label) if part_label else ''}"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(40, 53, 147)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.0), Inches(2.0), Inches(7.5), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(200, 200, 200)
    line.line.fill.background()

    content_left = Inches(2.0)
    content_top = Inches(2.4)
    content_width = Inches(7.5)
    content_height = Inches(4.5)

    if section_type in ['red_flags', 'security_risks']:
        warning_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            content_left - Inches(0.1),
            content_top - Inches(0.1),
            content_width + Inches(0.2),
            content_height + Inches(0.2)
        )
        if section_type == 'red_flags':
            bg_color = RGBColor(255, 243, 224)
            border_color = RGBColor(230, 81, 0)
        else:
            bg_color = RGBColor(255, 235, 238)
            border_color = RGBColor(198, 40, 40)
        warning_box.fill.solid()
        warning_box.fill.fore_color.rgb = bg_color
        warning_box.line.color.rgb = border_color
        warning_box.line.width = Pt(2)

    content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.TOP

    if isinstance(content_points, list):
        for i, point in enumerate(content_points):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            trimmed = point.strip()
            p.text = f"• {trimmed[:350]}"
            p.font.size = Pt(14)
            p.font.name = "Calibri"
            p.space_after = Pt(8)
            if section_type in ['red_flags', 'security_risks']:
                p.font.color.rgb = RGBColor(198, 40, 40)
            else:
                p.font.color.rgb = RGBColor(33, 33, 33)
    else:
        text = content_points if isinstance(content_points, str) else '\n\n'.join(content_points)
        p = text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(14)
        p.font.name = "Calibri"
        p.line_spacing = 1.2
        p.font.color.rgb = RGBColor(33, 33, 33)
    return slide


def add_table_slide_with_image(prs, title, emoji, items, section_type, part_label=None):
    """Add table slide for tech stack"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_decorative_element(slide, section_type)
    add_icon_shape(slide, emoji, 0.7, 1.5, 1.0)

    title_box = slide.shapes.add_textbox(Inches(2.0), Inches(1.3), Inches(7.5), Inches(0.6))
    p = title_box.text_frame.paragraphs[0]
    p.text = f"{title}{(' — ' + part_label) if part_label else ''}"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(40, 53, 147)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.0), Inches(2.0), Inches(7.5), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(200, 200, 200)
    line.line.fill.background()

    left = Inches(2.0)
    top = Inches(2.4)
    width = Inches(7.5)
    height = Inches(4.5)

    rows = min(len(items) + 1, 11)
    cols = 1
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    table.columns[0].width = width

    # Header
    cell = table.cell(0, 0)
    cell.text = "Technology / Framework / Library"
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(232, 234, 246)
    paragraph = cell.text_frame.paragraphs[0]
    paragraph.font.bold = True
    paragraph.font.size = Pt(14)
    paragraph.alignment = PP_ALIGN.CENTER

    # Data rows
    for i, item in enumerate(items[:10], start=1):
        cell = table.cell(i, 0)
        cell.text = f"  •  {item[:260]}"
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.size = Pt(12)
        paragraph.alignment = PP_ALIGN.LEFT
        if i % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(250, 250, 250)

    return slide


def generate_ppt(report_text, project_name, output_folder="analysis_output"):
    """Generate professional PowerPoint with auto-chunking"""
    print("\n" + "="*70)
    print("📊 GENERATING POWERPOINT PRESENTATION")
    print("="*70)
    
    os.makedirs(output_folder, exist_ok=True)
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    ppt_name = f"presentation_{timestamp}.pptx"
    output_path = os.path.join(output_folder, ppt_name)

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_ppt_title_slide(prs, project_name)

    sections = parse_analysis_report(report_text)
    if not sections:
        print("⚠️  Warning: No sections found in the report")
        return None

    section_config = {
        'project_summary': {'title': 'Project Summary', 'emoji': '📋', 'type': 'text'},
        'tech_stack': {'title': 'Technology Stack', 'emoji': '🛠️', 'type': 'table'},
        'file_structure': {'title': 'File & Folder Structure', 'emoji': '📁', 'type': 'bullet'},
        'core_modules': {'title': 'Core Modules', 'emoji': '⚙️', 'type': 'bullet'},
        'data_flow': {'title': 'Data Flow & Logic', 'emoji': '🔄', 'type': 'bullet'},
        'code_quality': {'title': 'Code Quality Analysis', 'emoji': '✅', 'type': 'bullet'},
        'red_flags': {'title': 'Red Flags & Issues', 'emoji': '⚠️', 'type': 'red_flags'},
        'security_risks': {'title': 'Security Risks', 'emoji': '🔒', 'type': 'security_risks'},
        'performance_risks': {'title': 'Performance Risks', 'emoji': '⚡', 'type': 'longtext'},
        'refactor_suggestions': {'title': 'Refactor Suggestions', 'emoji': '💡', 'type': 'bullet'}
    }

    # Chunking parameters
    TECH_ITEMS_PER_SLIDE = 10
    BULLETS_PER_SLIDE = 6
    TEXT_CHARS_PER_SLIDE = 800

    for section_key, cfg in section_config.items():
        if section_key not in sections or not sections[section_key].strip():
            continue
        
        content = sections[section_key]
        title = cfg['title']
        emoji = cfg['emoji']
        stype = cfg['type']

        print(f"📄 Creating slide(s): {title}")

        if stype == 'table' and section_key == 'tech_stack':
            tech_items = extract_bullet_points(content)
            if not tech_items:
                continue
            for i in range(0, len(tech_items), TECH_ITEMS_PER_SLIDE):
                part_items = tech_items[i:i + TECH_ITEMS_PER_SLIDE]
                part_label = f"Part {i // TECH_ITEMS_PER_SLIDE + 1}" if len(tech_items) > TECH_ITEMS_PER_SLIDE else None
                add_table_slide_with_image(prs, title, emoji, part_items, section_key, part_label=part_label)

        elif stype in ['bullet', 'red_flags', 'security_risks']:
            points = extract_bullet_points(content)
            if not points:
                continue
            for i in range(0, len(points), BULLETS_PER_SLIDE):
                part_points = points[i:i + BULLETS_PER_SLIDE]
                part_label = f"Part {i // BULLETS_PER_SLIDE + 1}" if len(points) > BULLETS_PER_SLIDE else None
                add_content_slide_with_image(prs, title, emoji, part_points, cfg['type'], part_label=part_label)

        elif stype == 'longtext':
            text = content.strip()
            if len(text) <= TEXT_CHARS_PER_SLIDE:
                add_content_slide_with_image(prs, title, emoji, text, 'text')
            else:
                sentences = re.split(r'(?<=[.!?])\s+', text)
                cur = ""
                part = 1
                for s in sentences:
                    if len(cur) + len(s) + 1 <= TEXT_CHARS_PER_SLIDE:
                        cur += (s + " ")
                    else:
                        part_label = f"Part {part}"
                        add_content_slide_with_image(prs, title, emoji, cur.strip(), 'text', part_label=part_label)
                        part += 1
                        cur = s + " "
                if cur.strip():
                    part_label = f"Part {part}" if part > 1 else None
                    add_content_slide_with_image(prs, title, emoji, cur.strip(), 'text', part_label=part_label)

        elif stype == 'text':
            text = content.strip()
            if len(text) <= TEXT_CHARS_PER_SLIDE:
                add_content_slide_with_image(prs, title, emoji, text, 'text')
            else:
                sentences = re.split(r'(?<=[.!?])\s+', text)
                cur = ""
                part = 1
                for s in sentences:
                    if len(cur) + len(s) + 1 <= TEXT_CHARS_PER_SLIDE:
                        cur += (s + " ")
                    else:
                        part_label = f"Part {part}"
                        add_content_slide_with_image(prs, title, emoji, cur.strip(), 'text', part_label=part_label)
                        part += 1
                        cur = s + " "
                if cur.strip():
                    part_label = f"Part {part}" if part > 1 else None
                    add_content_slide_with_image(prs, title, emoji, cur.strip(), 'text', part_label=part_label)

    # Closing slide
    print("🎬 Creating closing slide...")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(26, 35, 126)
    bg.line.fill.background()
    
    thank_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(6), Inches(2))
    text_frame = thank_box.text_frame
    
    p = text_frame.paragraphs[0]
    p.text = "✨ Thank You!"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    p2 = text_frame.add_paragraph()
    p2.text = "End of Analysis Report"
    p2.font.size = Pt(24)
    p2.font.color.rgb = RGBColor(200, 200, 200)
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(20)

    try:
        prs.save(output_path)
        print(f"✅ PowerPoint created successfully: {output_path}\n")
        return output_path
    except Exception as e:
        print(f"❌ Error saving PowerPoint: {e}")
        raise


# =============================================================================
# MAIN ORCHESTRATOR
# =============================================================================

def main():
    """Main orchestrator function"""
    print("\n" + "="*70)
    print("🚀 INTEGRATED GITHUB REPOSITORY ANALYZER")
    print("="*70)
    print("Features: Download → Analyze → Generate PDF & PPT")
    print("="*70 + "\n")
    
    # Step 1: Get repository URL or folder path
    choice = input("Enter '1' for GitHub URL or '2' for local folder: ").strip()
    
    if choice == '1':
        repo_url = input("📎 Enter GitHub repository URL: ").strip()
        if not repo_url:
            print("❌ No URL provided")
            return
        
        folder_path = download_github_repo(repo_url)
        if not folder_path:
            print("❌ Failed to download repository")
            return
        
        project_name = repo_url.split('/')[-1].replace('.git', '')
    
    elif choice == '2':
        folder_path = input("📂 Enter local folder path: ").strip()
        if not folder_path or not os.path.exists(folder_path):
            print("❌ Invalid folder path")
            return
        project_name = os.path.basename(os.path.abspath(folder_path))
    
    else:
        print("❌ Invalid choice")
        return
    
    # Step 2: Analyze repository
    max_files = input(f"📊 Max files to analyze (default=40): ").strip()
    max_files = int(max_files) if max_files.isdigit() else 40
    
    result = analyze_repository(folder_path, max_files, project_name)
    if not result:
        print("❌ Analysis failed")
        return
    
    report_text, report_path, project_name = result
    
    # Step 3: Generate documents
    gen_docs = input("\n📄 Generate PDF and PPT? (Y/n): ").strip().lower()
    if gen_docs != 'n':
        pdf_path = generate_pdf(report_text, project_name)
        ppt_path = generate_ppt(report_text, project_name)
        
        print("\n" + "="*70)
        print("🎉 SUCCESS! ALL DOCUMENTS GENERATED")
        print("="*70)
        print(f"📄 Text Report: {os.path.abspath(report_path)}")
        print(f"📄 PDF Report: {os.path.abspath(pdf_path)}")
        print(f"📊 PowerPoint: {os.path.abspath(ppt_path)}")
        print(f"📂 Output Folder: {os.path.abspath('analysis_output')}")
        print("="*70 + "\n")
    else:
        print(f"\n✅ Analysis complete! Report saved to: {report_path}\n")


if __name__ == "__main__":
    try:
        main()
    except KeyboardInterrupt:
        print("\n\n⚠️  Process interrupted by user")
    except Exception as e:
        print(f"\n❌ Error: {e}")
        import traceback
        traceback.print_exc()